# -*- coding: utf-8 -*-
"""生成《e1笔记-ANSYS学习总结.pptx》—— 杂志编辑风，第 1-11 页"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from deck_lib import *
from deck_lib import _line
from pptx import Presentation
from pptx.util import Emu

OUT = r"C:\obsidian\我的仓库\e1笔记-ANSYS学习总结.pptx"

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]


def new_slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(BLACKPG if dark else PAPER)
    return s


def subhead(s, x, y, txt, size=13.5, color=INK, w=5.0):
    return text(s, x, y, w, 0.34, [[(txt, size, True, color)]])


def mono_num(s, x, y, n, size=13, color=RED, w=0.6):
    return text(s, x, y, w, 0.32, [[("%02d" % n, size, True, color, F_MONO)]])


CH = {  # 章节运行头
    1: ("SEC.01 — GEOMETRY", "几何建模"), 2: ("SEC.02 — MESHING", "网格划分"),
    3: ("SEC.03 — MECHANICAL", "Mechanical"), 4: ("SEC.04 — LINEAR", "线性分析"),
    5: ("SEC.05 — DYNAMICS", "结构动力学"), 6: ("SEC.06 — NONLINEAR", "非线性分析"),
    7: ("SEC.07 — THERMAL", "热力学分析"), 8: ("SEC.08 — FATIGUE & MORE", "疲劳与其他"),
}


def content_page(page, chap, kick, title, standfirst=None):
    s = new_slide()
    en, cn = CH[chap]
    page_frame(s, en, cn, page)
    ghost(s, 10.55, 0.92, "%02d" % chap, size=88, w=1.9, h=1.5)
    title_block(s, kick, title, standfirst)
    return s


# ================= S1 封面 =================
s = new_slide()
text(s, 0.9, 0.72, 8.0, 0.3,
     [[("VOL.01 — ANSYS WORKBENCH · FINITE ELEMENT ANALYSIS NOTES", 10, True, INK, F_MONO, 180)]])
text(s, 9.4, 0.72, 3.03, 0.3, [[("2026 / 09", 10, True, GRAY, F_MONO, 180)]], align="r")
rule(s, 0.9, 1.22, 11.53, color=INK, weight=1.2)
ghost(s, 8.55, 1.55, "FE", size=150, w=3.9, h=2.6)
text(s, 0.86, 2.15, 10.5, 1.3, [[("有限元分析", 64, True, INK)]])
text(s, 0.86, 3.42, 10.5, 1.3,
     [[("学习", 64, True, INK), ("总结", 64, True, RED)]])
text(s, 0.9, 5.05, 9.6, 0.75,
     [{"r": [("几何建模 · 网格划分 · Mechanical · 线性 / 非线性 / 动力学 / 热 · 疲劳与多物理场", 13, False, GRAY)],
       "o": {"line": 1.5}},
      {"r": [("基于 e1-2 ~ e1-10 课程学习笔记的系统整理", 11, False, FAINT, F_MONO)],
       "o": {"line": 1.5, "sb": 6}}])
rule(s, 0.9, 6.45, 11.53, color=INK, weight=1.2)
idx = [("01", "建模"), ("02", "网格"), ("03", "MECHANICAL"), ("04", "线性"), ("05", "动力学"),
       ("06", "非线性"), ("07", "热学"), ("08", "疲劳·多场")]
runs = []
for i, (n, t1) in enumerate(idx):
    runs.append((n + " ", 9.5, True, RED, F_MONO))
    runs.append((t1 + ("　　" if i < 7 else ""), 9.5, False, INK, F_MONO))
text(s, 0.9, 6.62, 11.53, 0.3, [runs])

# ================= S2 目录 =================
s = new_slide()
page_frame(s, "CONTENTS", "目录", 2)
kicker(s, 0.9, 1.05, "CONTENTS — 本期内容")
text(s, 0.9, 1.36, 8.0, 0.75, [[("目录", 32, True, INK)]])
toc = [
    ("01", "几何建模", "DM 五要点 · SC 定位", "P3"),
    ("02", "网格划分", "五类网格 · 划分方法 · 国标", "P4"),
    ("03", "Mechanical", "三阶段流程 · 材料库", "P6"),
    ("04", "线性分析", "前提假设 · 静力与六类动力学", "P8"),
    ("05", "结构动力学", "模态 / 谱 / 谐响应 / 屈曲 / 瞬态", "P10"),
    ("06", "非线性分析", "类型 · 算法 · 单元与大变形", "P15"),
    ("07", "热力学分析", "传热方式 · 分析流程", "P17"),
    ("08", "疲劳与其他", "疲劳 · 刚体动力学 · LS-DYNA · Fluent", "P19"),
]
for i, (n, t1, t2, pg) in enumerate(toc):
    y = 2.42 + i * 0.55
    text(s, 0.9, y + 0.04, 0.6, 0.3, [[(n, 12, True, RED, F_MONO)]])
    text(s, 1.62, y, 2.9, 0.34, [[(t1, 15, True, INK)]])
    text(s, 4.7, y + 0.05, 6.0, 0.3, [[(t2, 10.5, False, GRAY)]])
    text(s, 10.8, y + 0.04, 1.63, 0.3, [[(pg, 11, True, INK, F_MONO)]], align="r")
    if i < 7:
        rule(s, 0.9, y + 0.47, 11.53)

# ================= S3 e1-2 几何建模 =================
s = content_page(3, 1, "SECTION 01 · GEOMETRY — e1-2", "DesignModeler 与 SpaceClaim",
                 "DM 的五大建模要点，构成本章主线。")
rows = [
    ("几何创建", "2D 草图绘制，3D 实体 / 面体 / 线体建模"),
    ("几何导入 / 导出", "STEP、IGES、Parasolid、CATIA 等主流格式"),
    ("建模操作", "拉伸、旋转、扫掠、放样、布尔运算、倒角、抽壳"),
    ("参数化设计", "尺寸 / 约束设为参数，联动 Workbench 优化模块"),
    ("冻结 / 非冻结", "控制特征是否参与布尔运算，是多体建模的基础"),
]
for i, (t1, t2) in enumerate(rows):
    y = 2.72 + i * 0.82
    mono_num(s, 0.9, y + 0.03, i + 1)
    text(s, 1.62, y, 6.3, 0.34, [[(t1, 14, True, INK)]])
    text(s, 1.62, y + 0.36, 6.3, 0.3, [[(t2, 11.5, False, GRAY)]])
    if i < 4:
        rule(s, 1.62, y + 0.7, 6.1)
vline(s, 8.5, 2.72, 6.55)
kicker(s, 8.85, 2.72, "SIDE NOTE — SC", size=9.5)
text(s, 8.85, 3.08, 3.58, 0.36, [[("SpaceClaim（SC）", 14.5, True, INK)]])
text(s, 8.85, 3.52, 3.58, 1.1,
     [{"r": [("直接建模思路的三维 CAD 系统。本课程以 UG 建模为主，未对 SC 做过多介绍。", 11.5, False, GRAY)],
       "o": {"line": 1.45}}])
rule(s, 8.85, 4.75, 3.58)
kicker(s, 8.85, 4.95, "KEYWORDS", size=9.5)
text(s, 8.85, 5.3, 3.58, 0.8,
     [{"r": [("草图 · 实体 · 布尔", 13, True, INK)], "o": {"line": 1.4}},
      {"r": [("参数化 · 冻结 · 多体", 13, True, INK)], "o": {"line": 1.4}}])

# ================= S4 e1-3 网格① =================
s = content_page(4, 2, "SECTION 02 · MESHING — e1-3", "五类网格与参考标准",
                 "Workbench 针对不同物理场，提供相应的网格划分器。")
types = [
    ("Meshing", "结构与热", "结构及热力学有限元网格（最常用）"),
    ("Electromagnetics", "电磁", "电磁场分析专用网格"),
    ("CFD", "计算流体动力学", "流体力学分析网格"),
    ("Explicit", "显式动力学", "显式动力学分析网格"),
    ("水动力学", "海洋 / 船舶", "常应用于海洋装备、船舶领域"),
]
for i, (en, cn, desc) in enumerate(types):
    y = 2.72 + i * 0.62
    text(s, 0.9, y + 0.02, 2.5, 0.3, [[(en, 12.5, True, INK, F_MONO)]])
    text(s, 3.45, y + 0.02, 1.75, 0.3, [[(cn, 11.5, True, RED)]])
    text(s, 5.25, y + 0.03, 3.1, 0.3, [[(desc, 10.5, False, GRAY)]])
    if i < 4:
        rule(s, 0.9, y + 0.46, 7.45)
sq(s, 0.9, 6.13, 0.1, color=RED)
text(s, 1.12, 6.04, 8.0, 0.34,
     [[("参考标准　", 12.5, True, INK),
       ("GB/T 33582-2017", 12.5, True, RED, F_MONO),
       ("《机械产品结构有限元力学分析通用规则》", 12.5, False, INK)]])
vline(s, 8.75, 2.72, 6.5)
gx0, gy0, st = 9.35, 3.0, 0.4
for i in range(4):
    _line(s, gx0 + i * st, gy0, gx0 + i * st, gy0 + 2 * st, GRAY, 1.1)
for j in range(3):
    _line(s, gx0, gy0 + j * st, gx0 + 3 * st, gy0 + j * st, GRAY, 1.1)
text(s, 9.15, gy0 + 2 * st + 0.08, 1.8, 0.26, [[("QUAD / HEX", 9, False, GRAY, F_MONO)]], align="c")
tx0 = 10.85
for i in range(4):
    _line(s, tx0 + i * st, gy0, tx0 + i * st, gy0 + 2 * st, GRAY, 1.1)
for j in range(3):
    _line(s, tx0, gy0 + j * st, tx0 + 3 * st, gy0 + j * st, GRAY, 1.1)
for j in range(2):
    for i in range(3):
        _line(s, tx0 + i * st, gy0 + j * st, tx0 + (i + 1) * st, gy0 + (j + 1) * st, RED, 1.1)
text(s, 10.65, gy0 + 2 * st + 0.08, 1.8, 0.26, [[("TRI / TET", 9, False, GRAY, F_MONO)]], align="c")
text(s, 9.05, 4.75, 3.45, 0.9,
     [{"r": [("结构网格规则、节点排布有序；非结构网格适应复杂几何，节点排布灵活。", 10.5, False, GRAY)],
       "o": {"line": 1.4}}])

# ================= S5 e1-3 网格② =================
s = content_page(5, 2, "SECTION 02 · MESHING — e1-3", "网格划分方法")
text(s, 0.9, 2.5, 5.4, 0.36,
     [[("3D ", 15, True, RED, F_MONO), ("三维划分", 14.5, True, INK),
       ("　五种方法", 10.5, False, FAINT)]])
m3 = [
    ("自动网格划分", "默认方法，快速获得整体网格"),
    ("四面体网格划分", "适应性最强，适合任意复杂几何"),
    ("六面体主导划分", "尽量生成六面体，过渡区域四面体"),
    ("扫略法", "源面网格沿路径扫略，规则高效"),
    ("多区法", "自动分块，块内结构化划分"),
]
for i, (t1, t2) in enumerate(m3):
    y = 3.06 + i * 0.72
    text(s, 0.9, y + 0.02, 0.45, 0.3, [[("%d" % (i + 1), 12, True, RED, F_MONO)]])
    text(s, 1.4, y, 4.9, 0.32, [[(t1, 13, True, INK)]])
    text(s, 1.4, y + 0.34, 4.9, 0.3, [[(t2, 10.5, False, GRAY)]])
vline(s, 6.6, 2.55, 6.55)
text(s, 7.05, 2.5, 5.4, 0.36,
     [[("2D ", 15, True, RED, F_MONO), ("二维划分", 14.5, True, INK),
       ("　三种方法", 10.5, False, FAINT)]])
m2 = ["四边形主导网格划分", "三角形主导网格划分", "四边形 / 三角形主导网格划分"]
for i, t1 in enumerate(m2):
    y = 3.06 + i * 0.52
    text(s, 7.05, y + 0.02, 0.45, 0.3, [[("%d" % (i + 1), 12, True, RED, F_MONO)]])
    text(s, 7.55, y, 4.9, 0.32, [[(t1, 13, True, INK)]])
rule(s, 7.05, 4.85, 5.38)
text(s, 7.05, 5.05, 5.38, 0.34, [[("划分思路", 13, True, INK)]])
text(s, 7.05, 5.45, 5.38, 1.0,
     [{"r": [("规则几何优先结构化方法（六面体 / 扫略 / 多区）；复杂几何用四面体或自动划分。", 11.5, False, GRAY)],
       "o": {"line": 1.45}}])

# ================= S6 e1-4 Mechanical① =================
s = content_page(6, 3, "SECTION 03 · MECHANICAL — e1-4", "集成式求解器与三阶段流程",
                 "求解能力：静力学 / 动力学 / 线性与非线性结构 / 热力学 / 磁场优化")
stages = [
    ("1", "前处理", ["几何模型系统的构建", "材料模型系统的构建", "有限元系统模型的构建"]),
    ("2", "求解", ["载荷边界条件", "位移边界条件", "求解设定"]),
    ("3", "后处理", ["结果趋势判定", "结果量级判定", "结果误差分析"]),
]
for i, (num, t1, items) in enumerate(stages):
    x = 0.9 + i * 3.95
    text(s, x, 2.6, 1.6, 0.85, [[(num, 40, True, GHOST, F_MONO)]])
    text(s, x, 3.42, 3.4, 0.4, [[(t1, 16.5, True, INK)]])
    rule(s, x, 3.94, 3.4, color=INK, weight=1.0)
    bullet_list(s, x, 4.14, 3.5, items, size=11.5, gap=8, color=RED, tcolor=GRAY, h=1.9)
    if i < 2:
        vline(s, x + 3.68, 2.75, 6.3)

# ================= S7 e1-4 Mechanical② =================
s = content_page(7, 3, "SECTION 03 · MECHANICAL — e1-4", "材料库与材料定义")
text(s, 0.9, 2.6, 5.35, 0.36, [[("工程数据源", 14.5, True, INK)]])
text(s, 0.9, 3.04, 5.35, 0.8,
     [{"r": [("直接调用内置工程数据（Engineering Data），为各类分析提供材料参数", 11.5, False, GRAY)],
       "o": {"line": 1.45}}])
rule(s, 0.9, 4.05, 5.35)
text(s, 0.9, 4.28, 5.35, 0.36, [[("自定义新材料", 14.5, True, INK)]])
text(s, 0.9, 4.72, 5.35, 1.0,
     [{"r": [("可新建材料并定义属性：", 11.5, False, GRAY),
             ("各向同性弹性（杨氏模量、泊松比）、比热容等", 11.5, True, INK)],
       "o": {"line": 1.45}}])
vline(s, 6.5, 2.6, 6.1)
text(s, 6.9, 2.6, 5.5, 0.36, [[("常用材料参数", 14.5, True, INK)]])
props = [
    ("E", "弹性模量", "刚度基本量，决定抵抗变形的能力"),
    ("ν", "泊松比", "横向应变与轴向应变之比"),
    ("ρ", "密度", "模态、重力、惯性计算需要"),
    ("c", "比热容", "瞬态热分析升温计算需要"),
]
for i, (sym, t1, t2) in enumerate(props):
    y = 3.14 + i * 0.66
    text(s, 6.9, y, 0.5, 0.32, [[(sym, 14, True, RED, F_MONO)]])
    text(s, 7.5, y + 0.03, 1.4, 0.3, [[(t1, 12.5, True, INK)]])
    text(s, 8.95, y + 0.05, 3.45, 0.3, [[(t2, 10.5, False, GRAY)]])
    if i < 3:
        rule(s, 6.9, y + 0.5, 5.5)
sq(s, 0.9, 6.32, 0.1, color=RED)
text(s, 1.12, 6.23, 11.3, 0.34,
     [[("模态等动力学分析的材料定义三要素：", 12.5, False, INK),
       ("弹性模量 E + 泊松比 ν + 密度 ρ", 12.5, True, RED)]])

# ================= S8 e1-5 线性分析① =================
s = content_page(8, 4, "SECTION 04 · LINEAR — e1-5", "定义、前提与弹性力学假设")
text(s, 0.9, 2.52, 11.53, 0.75,
     [{"r": [("线性分析是最基本、应用最广泛的一类分析 —— 适用于", 13.5, False, INK),
             ("线弹性材料、静态或动态稳定状态加载", 13.5, True, RED),
             ("的工况。", 13.5, False, INK)], "o": {"line": 1.45}}])
rule(s, 0.9, 3.42, 11.53, color=INK, weight=1.0)
subhead(s, 0.9, 3.62, "线性分析前提")
prem = [
    ("材料线性", "应力与应变关系呈线性状态"),
    ("小位移 · 小应变 · 小转动", "几何上处于小变形范围"),
    ("刚度不变", "刚度不随结构变形发生变化"),
]
for i, (t1, t2) in enumerate(prem):
    y = 4.12 + i * 0.78
    mono_num(s, 0.9, y + 0.02, i + 1)
    text(s, 1.62, y, 4.6, 0.32, [[(t1, 13, True, INK)]])
    text(s, 1.62, y + 0.35, 4.6, 0.3, [[(t2, 11, False, GRAY)]])
vline(s, 6.6, 3.62, 6.55)
subhead(s, 6.98, 3.62, "弹性力学五大假设")
asm = [
    ("连续性", "物质无间隙地充满所在空间"),
    ("线弹性", "应力与应变成正比，卸载完全恢复"),
    ("均匀性", "各点材料属性相同"),
    ("各向同性", "各方向材料属性相同"),
    ("小变形", "变形远小于结构特征尺寸"),
]
for i, (t1, t2) in enumerate(asm):
    y = 4.14 + i * 0.5
    sq(s, 6.98, y + 0.09, 0.09, color=RED)
    text(s, 7.2, y, 1.4, 0.3, [[(t1, 12.5, True, INK)]])
    text(s, 8.62, y + 0.02, 3.8, 0.3, [[(t2, 10.5, False, GRAY)]])

# ================= S9 e1-5 线性分析② =================
s = content_page(9, 4, "SECTION 04 · LINEAR — e1-5", "线性分析内容体系")
text(s, 0.9, 2.6, 4.7, 0.4, [[("线性静力学分析", 15, True, INK)]])
text(s, 0.9, 3.1, 4.7, 1.0,
     [{"r": [("系统运动速度为 0，分析平衡状态下结构的受力与变形 —— 最基础的分析类型。", 11.5, False, GRAY)],
       "o": {"line": 1.45}}])
text(s, 0.9, 4.45, 4.7, 0.9,
     [[("K X = ", 42, True, INK, F_MONO), ("F", 42, True, RED, F_MONO)]], align="c")
text(s, 0.9, 5.5, 4.7, 0.3, [[("刚度矩阵 × 位移 = 载荷", 10, False, FAINT, F_MONO)]], align="c")
vline(s, 6.0, 2.6, 6.55)
text(s, 6.4, 2.6, 6.0, 0.4,
     [[("线性动力学分析 ", 15, True, INK), ("—— 六种类型", 12, False, FAINT)]])
dyn = [
    ("模态分析", "求固有频率与振型"),
    ("谐响应分析", "正弦稳态受迫振动"),
    ("随机振动分析", "PSD 功率谱统计响应"),
    ("响应谱分析", "频域峰值响应"),
    ("瞬态动力学", "载荷随时间变化"),
    ("线性屈曲分析", "临界载荷与稳定性"),
]
for i, (t1, t2) in enumerate(dyn):
    x = 6.4 + (i % 2) * 3.1
    y = 3.22 + (i // 2) * 1.08
    mono_num(s, x, y, i + 1, size=11.5)
    text(s, x + 0.52, y - 0.02, 2.6, 0.32, [[(t1, 13, True, INK)]])
    text(s, x + 0.52, y + 0.32, 2.6, 0.3, [[(t2, 10.5, False, GRAY)]])
    if i < 4:
        rule(s, x, y + 0.78, 2.75)

# ================= S10 e1-6 动力学总览 =================
s = content_page(10, 5, "SECTION 05 · DYNAMICS — e1-6", "动力学分析总览")
vline(s, 6.65, 2.55, 6.55)
rule(s, 0.9, 4.52, 11.53)
text(s, 0.9, 2.62, 5.4, 0.3, [[("A", 10, True, RED, F_MONO), ("　求解类型", 13.5, True, INK)]])
text(s, 0.9, 3.1, 5.4, 0.5, [[("瞬态 · 冲击 · 碰撞", 20, True, INK)]])
text(s, 0.9, 3.75, 5.4, 0.35, [[("载荷随时间剧烈变化的动力问题", 11, False, GRAY)]])
text(s, 7.05, 2.62, 5.4, 0.3, [[("B", 10, True, RED, F_MONO), ("　求解方法", 13.5, True, INK)]])
text(s, 7.05, 3.12, 5.4, 0.36,
     [[("显式动力学", 13.5, True, INK), ("　高频冲击类", 11.5, False, GRAY)]])
text(s, 7.05, 3.6, 5.4, 0.36,
     [[("隐式动力学", 13.5, True, INK), ("　低频振动类", 11.5, False, GRAY)]])
text(s, 0.9, 4.72, 5.4, 0.3, [[("C", 10, True, RED, F_MONO), ("　惯性力", 13.5, True, INK)]])
text(s, 0.9, 5.2, 5.4, 0.85,
     [{"r": [("动力学问题必须考虑惯性力 —— 求解时需", 12, False, INK),
             ("开启时间积分效应", 12, True, RED), ("。", 12, False, INK)],
       "o": {"line": 1.45}}])
text(s, 7.05, 4.72, 5.4, 0.3, [[("D", 10, True, RED, F_MONO), ("　阻尼", 13.5, True, INK)]])
text(s, 7.05, 5.12, 5.38, 0.66,
     [{"r": [("常规情况下的能量消散；常见形式：", 11, False, GRAY),
             ("瑞利阻尼（α + β + 阻尼比）", 11, True, INK)], "o": {"line": 1.4}}])
segs = [("<1", "小阻尼"), ("1–5", "显著"), ("5–10", "非常显著"), (">10", "大阻尼")]
barx, barw = 7.05, 5.3
rule(s, barx, 6.18, barw, color=INK, weight=1.2)
for i in range(5):
    dot(s, barx + barw * i / 4, 6.18, 0.075, color=RED)
for i, (v, lb) in enumerate(segs):
    cx = barx + barw * (i + 0.5) / 4
    text(s, cx - 0.66, 6.28, 1.32, 0.3,
         [[(v, 9.5, True, RED, F_MONO), (" " + lb, 9, False, GRAY)]], align="c")

# ================= S11 e1-6 模态分析 =================
s = content_page(11, 5, "SECTION 05 · DYNAMICS — e1-6", "模态分析")
subhead(s, 0.9, 2.55, "作用", size=14)
acts = [
    "确定结构自身的振动特性：固有频率与振型",
    "振型为无量纲量，可反映结构的相对强弱",
    "改进结构固有频率或工况，避免共振",
    "是响应谱、随机振动、谐响应等分析的基础",
]
bullet_list(s, 0.9, 3.02, 6.2, acts, size=12, gap=8, h=2.0)
rule(s, 0.9, 5.25, 6.2, color=INK, weight=1.0)
subhead(s, 0.9, 5.45, "种类", size=14)
text(s, 0.9, 5.9, 6.3, 0.36,
     [[("自由模态 ", 13, True, INK), ("/ ", 13, True, RED, F_MONO),
       ("约束模态 ", 13, True, INK), ("/ ", 13, True, RED, F_MONO),
       ("有预应力模态 ", 13, True, INK), ("/ ", 13, True, RED, F_MONO),
       ("含接触模态", 13, True, INK)]])
vline(s, 7.6, 2.55, 6.55)
subhead(s, 8.0, 2.55, "实质", size=14)
text(s, 8.0, 3.0, 4.4, 0.75,
     [{"r": [("求解方程的", 11.5, False, GRAY), ("特征值 λ", 11.5, True, INK),
             ("（自振圆频率的平方）与", 11.5, False, GRAY), ("特征向量 v", 11.5, True, INK),
             ("（振型）", 11.5, False, GRAY)], "o": {"line": 1.4}}])
text(s, 8.0, 3.85, 4.43, 0.7,
     [[("f = ", 30, True, INK, F_MONO), ("√λ", 30, True, RED, F_MONO),
       (" / 2π", 30, True, INK, F_MONO)]], align="c")
rule(s, 8.0, 4.85, 4.43, color=INK, weight=1.0)
subhead(s, 8.0, 5.05, "关键结论", size=14)
text(s, 8.0, 5.5, 4.43, 0.65,
     [{"r": [("模态频率仅与结构的", 11.5, False, GRAY), ("质量矩阵、刚度矩阵", 11.5, True, INK),
             ("有关", 11.5, False, GRAY)], "o": {"line": 1.4}}])
text(s, 8.0, 6.2, 4.43, 0.34,
     [[("材料三要素　", 11, False, GRAY),
       ("E + ν + ρ", 14, True, RED, F_MONO)]])

# ================= S12 e1-6 响应谱与随机振动 =================
s = content_page(12, 5, "SECTION 05 · DYNAMICS — e1-6", "响应谱与随机振动")
kicker(s, 0.9, 2.5, "RESPONSE SPECTRUM", size=9.5)
text(s, 0.9, 2.82, 5.5, 0.4, [[("响应谱分析", 15.5, True, INK)]])
label_rows(s, 0.9, 3.36, 5.5, [
    ("是什么", "利用频域分析技术计算结构的峰值响应"),
    ("输入", "频谱可为位移谱 / 速度谱 / 加速度谱"),
    ("输出", "各阶振型在给定输入下的最大响应，按“响应系数 × 振型”叠加为总体响应"),
    ("作用", "替代瞬态分析快速获取峰值响应，常用于地震响应谱分析"),
    ("技术", "单点谱分析 / 多点谱分析 / 阻尼比"),
    ("注意", "缺失质量效应、缺失刚度响应效应"),
], gap=0.13, size=11.5)
vline(s, 6.75, 2.55, 6.55)
kicker(s, 7.15, 2.5, "RANDOM VIBRATION", size=9.5)
text(s, 7.15, 2.82, 5.3, 0.4, [[("随机振动分析", 15.5, True, INK)]])
label_rows(s, 7.15, 3.36, 5.28, [
    ("是什么", "谱分析技术的特定情况 —— 功率谱密度（PSD）分析"),
    ("输入", "功率谱密度函数 PSD：随机载荷时间历程的统计响应"),
    ("输出", "响应的概率统计结果"),
    ("场景", "汽车路谱平整度、火箭发射等载荷不确定工况"),
    ("怎么做", "分析流程与响应谱技术类似"),
], gap=0.17, size=11.5)

# ================= S13 e1-6 谐响应分析 =================
s = content_page(13, 5, "SECTION 05 · DYNAMICS — e1-6", "谐响应分析",
                 "频率响应分析 / 扫频分析 —— 时域计算 + 稳态受迫振动")
text(s, 0.9, 2.72, 11.53, 0.55,
     [[("m·a + c·v + k·x = f(t) ,  f = H·cos(ωt + φ)", 21, True, INK, F_MONO)]], align="c")
rule(s, 0.9, 3.5, 11.53)
text(s, 0.9, 3.7, 5.4, 0.36, [[("结果表达形式", 13.5, True, INK)]])
text(s, 0.9, 4.14, 5.4, 0.85,
     [{"r": [("方式一　幅值 + 相位角", 12, True, INK)], "o": {"sa": 3}},
      {"r": [("响应 = 幅值·e^(iφ)，φ = arctan(虚部 / 实部)", 11, False, GRAY)]}])
text(s, 0.9, 5.14, 5.4, 0.85,
     [{"r": [("方式二　实部 + 虚部", 12, True, INK)], "o": {"sa": 3}},
      {"r": [("响应 = 实部 + i·虚部", 11, False, GRAY)]}])
vline(s, 6.75, 3.66, 6.3)
text(s, 7.15, 3.7, 5.3, 0.36, [[("作用", 13.5, True, INK)]])
bullet_list(s, 7.15, 4.14, 5.28, [
    "确定结构在已知幅值与频率的正弦荷载下的稳态响应",
    "确定偏心质量旋转产生的离心力",
    "输出随频率变化的幅值 / 相位曲线",
], size=11.5, gap=8, tcolor=GRAY, h=1.6)
rule(s, 0.9, 6.02, 11.53)
text(s, 0.9, 6.2, 2.2, 0.32, [[("频带划分", 13, True, INK)]])
text(s, 3.0, 6.22, 3.3, 0.32,
     [[("1 倍频程 ", 10.5, True, INK), ("f2 = 2·f1，fc = √(f1·f2)", 10.5, False, GRAY, F_MONO)]])
text(s, 6.6, 6.22, 2.9, 0.32,
     [[("1/2 频带 ", 10.5, True, INK), ("f2 = √2·f1 ≈ 1.41·f1", 10.5, False, GRAY, F_MONO)]])
text(s, 9.8, 6.22, 2.7, 0.32,
     [[("1/3 频带 ", 10.5, True, INK), ("f2 = 2^(1/3)·f1 ≈ 1.26·f1", 10.5, False, GRAY, F_MONO)]])

# ================= S14 e1-6 屈曲与瞬态 =================
s = content_page(14, 5, "SECTION 05 · DYNAMICS — e1-6", "线性屈曲与瞬态动力学")
kicker(s, 0.9, 2.5, "LINEAR BUCKLING", size=9.5)
text(s, 0.9, 2.82, 5.5, 0.4, [[("线性屈曲分析", 15.5, True, INK)]])
label_rows(s, 0.9, 3.36, 5.5, [
    ("是什么", "结构在垂直轴向荷载下，因微小扰动出现系统崩溃 —— 即结构的稳定性计算"),
    ("作用", "确定临界载荷；相较非线性屈曲，可作为产品初始研发的评估指导手段"),
    ("常用对象", "细长杆、真空类零件"),
    ("怎么做", "需给定预应力条件（先施加预应力再求解特征值）"),
], gap=0.18, size=11.5)
vline(s, 6.75, 2.55, 6.55)
kicker(s, 7.15, 2.5, "TRANSIENT DYNAMICS", size=9.5)
text(s, 7.15, 2.82, 5.3, 0.4, [[("瞬态动力学分析", 15.5, True, INK)]])
label_rows(s, 7.15, 3.36, 5.28, [
    ("是什么", "载荷随时间变化的动力学分析"),
    ("作用", "确定任意时刻载荷作用下结构的受力状态（位移、应力、应变）"),
    ("怎么做", "通过时间步（时间步长）控制求解过程"),
    ("关联", "惯性力相关 —— 求解需开启时间积分效应"),
], gap=0.22, size=11.5)

# ================= S15 e1-7 非线性① =================
s = content_page(15, 6, "SECTION 06 · NONLINEAR — e1-7", "非线性行为、类型与求解算法")
text(s, 0.9, 2.52, 11.53, 0.45,
     [[("非线性本质：载荷引起", 13.5, False, INK), ("结构刚度发生变化", 13.5, True, RED),
       (" —— 线性问题中刚度矩阵 [k] 保持不变。", 13.5, False, INK)]])
rule(s, 0.9, 3.25, 11.53, color=INK, weight=1.0)
nls = [
    ("01", "几何非线性", "结构产生大变形导致明显形状变化，常见于细长杆与长薄板"),
    ("02", "材料非线性", "金属的弹塑性；橡胶材料的超弹性"),
    ("03", "状态非线性", "接触；生死单元"),
]
for i, (n, t1, t2) in enumerate(nls):
    x = 0.9 + i * 3.98
    text(s, x, 3.5, 1.0, 0.35, [[(n, 14, True, RED, F_MONO)]])
    text(s, x, 3.92, 3.5, 0.4, [[(t1, 15.5, True, INK)]])
    text(s, x, 4.4, 3.6, 0.95, [{"r": [(t2, 11.5, False, GRAY)], "o": {"line": 1.4}}])
    if i < 2:
        vline(s, x + 3.7, 3.55, 5.35)
rule(s, 0.9, 5.6, 11.53)
text(s, 0.9, 5.8, 5.4, 0.34, [[("求解算法", 13.5, True, INK)]])
text(s, 0.9, 6.22, 5.4, 0.35,
     [[("牛顿-拉夫逊迭代 · 多次平衡迭代 · 迭代步收敛", 11.5, False, GRAY)]])
text(s, 7.0, 5.8, 5.43, 0.34, [[("关键认识", 13.5, True, INK)]])
text(s, 7.0, 6.2, 5.43, 0.4,
     [[("收敛主要靠载荷步 / 子步控制；", 11.5, False, GRAY),
       ("结果收敛 ≠ 结果正确", 12.5, True, RED)]])

# ================= S16 e1-7 非线性② =================
s = content_page(16, 6, "SECTION 06 · NONLINEAR — e1-7", "单元控制与大变形分析")
text(s, 0.9, 2.55, 5.4, 0.36, [[("单元类型控制", 13.5, True, INK)]])
label_rows(s, 0.9, 3.0, 5.4, [
    ("隐式求解", "高阶单元 —— 适用于弯曲主导问题"),
    ("显式求解", "低阶单元 —— 注意剪切锁定、体积锁定"),
], gap=0.2, size=11.5)
rule(s, 0.9, 4.15, 5.4)
text(s, 0.9, 4.35, 5.4, 0.36, [[("单元求解控制", 13.5, True, INK)]])
text(s, 0.9, 4.8, 5.4, 0.7,
     [{"r": [("完全积分 ", 13, True, INK), ("/ ", 13, True, RED, F_MONO),
             ("缩减积分", 13, True, INK)], "o": {"sa": 4}},
      {"r": [("—— 不同的积分方案影响求解精度与沙漏控制", 10.5, False, GRAY)]}])
vline(s, 6.75, 2.55, 6.55)
text(s, 7.15, 2.55, 5.3, 0.36, [[("大变形分析", 13.5, True, INK)]])
text(s, 7.15, 2.95, 2.2, 0.8, [[("10%", 40, True, RED, F_MONO)]])
text(s, 7.15, 3.78, 5.28, 0.34,
     [[("经验罚则：横向位移超过厚度的 10% 时，需启用大变形分析", 10.5, False, GRAY)]])
label_rows(s, 7.15, 4.28, 5.28, [
    ("作用", "考虑大变形 / 大转动 / 大应变引起的单元形状与方向变化导致的刚度改变"),
    ("代价", "结果更精确，但需迭代求解、可能分步加载，计算时间长"),
    ("其他场景", "系统存在失稳时；使用超弹性材料时"),
], gap=0.13, size=11.5)
rule(s, 0.9, 6.02, 11.53)
text(s, 0.9, 6.22, 2.3, 0.32, [[("主要链接关系", 12.5, True, INK)]])
text(s, 3.1, 6.24, 9.3, 0.32,
     [[("布尔运算 · 共节点 · 运动副 · 接触 · 刚性耦合 · 柔性耦合 · 梁 / 杆 · CP · CE",
        11.5, False, GRAY)]])

# ================= S17 e1-8 热力学① =================
s = content_page(17, 7, "SECTION 07 · THERMAL — e1-8", "概念、作用与分析种类")
text(s, 0.9, 2.52, 11.53, 0.5,
     [[("研究宏观物质与", 13.5, False, INK), ("冷热有关的物理性质及其变化规律", 13.5, True, RED),
       ("的学科。", 13.5, False, INK)]])
rule(s, 0.9, 3.25, 11.53, color=INK, weight=1.0)
text(s, 0.9, 3.45, 5.4, 0.36, [[("作用", 13.5, True, INK)]])
bullet_list(s, 0.9, 3.9, 5.35, [
    "石油化工、核动力、制动器、压力容器的温度应力",
    "系统暂态过程中产生的瞬态温度应力",
], size=11.5, gap=8, tcolor=GRAY, h=1.1)
vline(s, 6.65, 3.45, 6.4)
text(s, 7.05, 3.45, 5.4, 0.36, [[("四种分析类型", 13.5, True, INK)]])
kinds = [
    ("a", "稳态热分析", "温度场不随时间变化"),
    ("b", "瞬态热分析", "温度场随时间变化"),
    ("c", "热结构间接耦合", "先热后结构，顺序求解"),
    ("d", "热结构直接耦合", "热-结构同时求解"),
]
for i, (n, t1, t2) in enumerate(kinds):
    y = 3.95 + i * 0.58
    text(s, 7.05, y + 0.02, 0.4, 0.3, [[(n, 11.5, True, RED, F_MONO)]])
    text(s, 7.5, y, 2.5, 0.3, [[(t1, 12.5, True, INK)]])
    text(s, 10.1, y + 0.02, 2.35, 0.3, [[(t2, 10, False, GRAY)]])
    if i < 3:
        rule(s, 7.05, y + 0.44, 5.38)
rule(s, 0.9, 6.02, 11.53)
text(s, 0.9, 6.25, 2.4, 0.34, [[("瞬态热力学一般方程", 12, True, INK)]])
text(s, 3.4, 6.16, 5.2, 0.45,
     [[("[C]·d{T}/dt + [K]·{T} = {Q}", 17, True, INK, F_MONO)]], align="c")
text(s, 8.9, 6.2, 3.53, 0.55,
     [{"r": [("[C] 比热 · [K] 传导 · {T} 节点温度向量", 9, False, FAINT, F_MONO)], "o": {"line": 1.35}},
      {"r": [("{Q} 热载荷向量", 9, False, FAINT, F_MONO)], "o": {"line": 1.35}}])

# ================= S18 e1-8 热力学② =================
s = content_page(18, 7, "SECTION 07 · THERMAL — e1-8", "三大传热方式与分析流程")
ht = [
    ("1", "热传导", "物体内部或系统内各部分存在温度差引起的热量传递"),
    ("2", "热对流", "温度不同的流体之间：高温附近空气膨胀 → 密度差 → 自然对流 / 强迫对流"),
    ("3", "热辐射", "物体发射电磁能进行热量交换；真空中更显著；黑体辐射为理想基准；多物体辐射高度非线性"),
]
for i, (n, t1, t2) in enumerate(ht):
    x = 0.9 + i * 3.98
    text(s, x, 2.55, 1.2, 0.8, [[(n, 36, True, GHOST, F_MONO)]])
    text(s, x, 3.35, 3.5, 0.4, [[(t1, 15.5, True, INK)]])
    text(s, x, 3.85, 3.6, 1.3, [{"r": [(t2, 11, False, GRAY)], "o": {"line": 1.4}}])
    if i < 2:
        vline(s, x + 3.7, 2.7, 5.0)
rule(s, 0.9, 5.35, 11.53, color=INK, weight=1.0)
text(s, 0.9, 5.55, 8.0, 0.36, [[("分析流程 —— 三大步三小步", 13.5, True, INK)]])
text(s, 0.9, 6.05, 11.53, 0.4,
     [[("前处理", 13, True, INK), ("（几何 · 材料 · 有限元构建）", 11.5, False, GRAY),
       ("  →  ", 13, True, RED, F_MONO),
       ("求解", 13, True, INK), ("（产热 · 散热 · 求解设定）", 11.5, False, GRAY),
       ("  →  ", 13, True, RED, F_MONO),
       ("后处理", 13, True, INK), ("（趋势 · 量级 · 误差）", 11.5, False, GRAY)]])

# ================= S19 e1-10 疲劳分析 =================
s = content_page(19, 8, "SECTION 08 · FATIGUE — e1-10", "疲劳分析")
text(s, 0.86, 2.45, 3.7, 1.55, [[("80%", 84, True, RED, F_MONO)]])
text(s, 0.9, 4.15, 3.6, 0.9,
     [{"r": [("机械设备的损坏源于疲劳失效 —— 疲劳是产品失效的主要根源。", 11.5, False, GRAY)],
       "o": {"line": 1.4}}])
rule(s, 0.9, 5.15, 3.6)
text(s, 0.9, 5.35, 3.6, 0.3, [[("高周疲劳", 13, True, INK), ("　应力疲劳", 10.5, False, GRAY)]])
text(s, 0.9, 5.75, 3.6, 0.3, [[("低周疲劳", 13, True, INK), ("　应变疲劳", 10.5, False, GRAY)]])
vline(s, 4.9, 2.5, 6.55)
text(s, 5.3, 2.55, 7.13, 0.36, [[("如何做 —— Workbench 疲劳工具（高周疲劳）", 13.5, True, INK)]])
label_rows(s, 5.3, 3.0, 7.1, [
    ("载荷", "恒定振幅载荷；成比例载荷 / 非比例载荷"),
    ("曲线", "应力-寿命（S-N）曲线 / 应变-寿命曲线"),
], gap=0.14, size=11.5)
rule(s, 5.3, 4.02, 7.13)
text(s, 5.3, 4.2, 7.13, 0.34, [[("S-N 曲线", 13.5, True, INK)]])
text(s, 5.3, 4.62, 7.13, 0.6,
     [{"r": [("试件疲劳测试所得结果；弯曲或轴向测试反映单轴应力状态；平均应力状态影响疲劳寿命。", 11.5, False, GRAY)],
       "o": {"line": 1.4}}])
text(s, 5.3, 5.42, 7.13, 0.34,
     [[("影响因素　", 11.5, True, INK),
       ("延展性 · 加工工艺 · 表面粗糙度 · 残余应力 · 载荷环境", 11.5, False, GRAY)]])
rule(s, 5.3, 5.98, 7.13)
sq(s, 5.3, 6.25, 0.1, color=RED)
text(s, 5.52, 6.16, 6.95, 0.65,
     [{"r": [("前提：线性静力分析，非线性问题慎用；接触用 ", 11.5, False, INK),
             ("bonded + no separation", 11.5, True, RED, F_MONO),
             (" 保证接触状态不随载荷波动变化。", 11.5, False, INK)], "o": {"line": 1.35}}])

# ================= S20 e1-10 三小节 =================
s = content_page(20, 8, "SECTION 08 · FATIGUE & MORE — e1-10", "刚体动力学 · LS-DYNA · Fluent")
cols = [
    ("RIGID BODY", "刚体动力学", [
        ("用途", "模拟各类运动机械及含弹簧系统的运动，各运动副的运动累积"),
        ("输入输出", "力、力矩、位移、速度、加速度"),
        ("特点", "阻尼可通过弹簧考虑；无应力应变结果"),
        ("实现", "显式求解法；材料仅需密度；几何用壳体 / 面体 / 实体"),
    ]),
    ("LS-DYNA", "LS-DYNA", [
        ("定位", "通用的非线性动力学分析软件"),
        ("能力", "爆炸、高速碰撞、金属成型等高度非线性冲击问题"),
        ("延伸", "还可求解传热、流体、振动等多场耦合问题"),
    ]),
    ("FLUENT CFD", "Fluent 流体力学", [
        ("定位", "对流体力学各类问题进行数值实验、计算机模拟与分析"),
        ("方法", "有限差分 FDM · 有限元 FEM · 有限体积 FVM"),
        ("模型", "定常与非定常、层流与湍流、可压与不可压流动、传热、化学反应等"),
    ]),
]
for i, (k, t1, rows) in enumerate(cols):
    x = 0.9 + i * 3.98
    kicker(s, x, 2.55, k, size=9.5)
    text(s, x, 2.88, 3.6, 0.4, [[(t1, 15.5, True, INK)]])
    label_rows(s, x, 3.42, 3.6, rows, gap=0.19, size=10.5, label_color=INK, text_color=GRAY)
    if i < 2:
        vline(s, x + 3.7, 2.6, 6.5)

# ================= S21 总结（深色收尾） =================
s = new_slide(dark=True)
text(s, 0.9, 0.7, 8.0, 0.3,
     [[("SUMMARY — 知识体系全景", 10, True, RED, F_MONO, 180)]])
text(s, 10.0, 0.7, 2.43, 0.3, [[("21 / 21", 10, False, DARKMUT, F_MONO, 150)]], align="r")
rule(s, 0.9, 1.18, 11.53, color=DARKRULE)
text(s, 0.9, 1.45, 11.0, 0.85, [[("有限元知识体系全景", 34, True, "FFFFFF")]])
text(s, 0.9, 2.5, 11.53, 0.4,
     [[("几何建模 ", 13.5, True, DARKTXT), ("→", 13.5, True, RED, F_MONO),
       (" 网格划分 ", 13.5, True, DARKTXT), ("→", 13.5, True, RED, F_MONO),
       (" Mechanical ", 13.5, True, DARKTXT), ("→", 13.5, True, RED, F_MONO),
       (" 线性分析 ", 13.5, True, DARKTXT), ("→", 13.5, True, RED, F_MONO),
       (" 结构动力学 ", 13.5, True, DARKTXT), ("→", 13.5, True, RED, F_MONO),
       (" 非线性 ", 13.5, True, DARKTXT), ("→", 13.5, True, RED, F_MONO),
       (" 热力学 ", 13.5, True, DARKTXT), ("→", 13.5, True, RED, F_MONO),
       (" 疲劳 · 多场", 13.5, True, DARKTXT)]])
rule(s, 0.9, 3.2, 11.53, color=DARKRULE)
takes = [
    ("01", "分析主线", [("前处理（几何 / 材料 / 有限元模型）→ 求解（边界与设定）→ 后处理（趋势 · 量级 · 误差）", False)]),
    ("02", "选型逻辑", [("先判断问题类型 —— 线性或非线性、静力或动力学、单场或耦合场，再选分析类型与算法", False)]),
    ("03", "结果心法", [("收敛 ≠ 正确", True), ("；从趋势、量级、误差三重校验，兼顾计算效率与精度", False)]),
]
for i, (n, t1, parts) in enumerate(takes):
    x = 0.9 + i * 3.98
    text(s, x, 3.55, 1.0, 0.4, [[(n, 16, True, RED, F_MONO)]])
    text(s, x, 4.05, 3.5, 0.4, [[(t1, 15.5, True, "FFFFFF")]])
    runs = [(t, 11.5, b, RED if b else DARKTXT) for t, b in parts]
    text(s, x, 4.55, 3.6, 1.7, [{"r": runs, "o": {"line": 1.5}}])
    if i < 2:
        vline(s, x + 3.7, 3.6, 6.3, color=DARKRULE)
rule(s, 0.9, 6.7, 11.53, color=DARKRULE)
text(s, 0.9, 6.85, 8.0, 0.3,
     [[("基于 e1-2 ~ e1-10 课程学习笔记整理 · 2026-09", 8.5, False, DARKMUT, F_MONO, 80)]])

# ================= 保存 =================
import shutil
prs.core_properties.title = "ANSYS Workbench 有限元分析学习总结"
prs.core_properties.author = "AI 整理"
OUT_WORK = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_out.pptx")
prs.save(OUT_WORK)
print("saved(work):", OUT_WORK)
try:
    shutil.copyfile(OUT_WORK, OUT)
    print("saved:", OUT)
except PermissionError:
    print("WARN: 正式文件被占用（可能正被打开），最新版暂存于 _out.pptx，关闭后重新运行即可覆盖")
