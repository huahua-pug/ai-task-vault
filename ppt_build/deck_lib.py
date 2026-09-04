# -*- coding: utf-8 -*-
"""PPT 生成公共库——杂志编辑风：白纸底、排版直接、细线分栏、单一红色强调"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# 16:9 画布
SW, SH = 13.3333, 7.5

# ---------- 配色：暖白纸 + 墨黑 + 编辑红 ----------
PAPER = "FDFCF9"   # 暖白纸底
INK = "141719"     # 近黑正文
GRAY = "6E7479"    # 弱化灰
FAINT = "9DA2A6"   # 更浅灰（页眉页脚）
RULE = "D9D5CB"    # 暖灰细线
RED = "C8102E"     # 编辑红（唯一强调色）
GHOST = "ECE9E0"   # 超大浅色序号
WHITE = "FFFFFF"
BLACKPG = "121417" # 总结页深色底
DARKTXT = "E9E7E2" # 深色页正文
DARKMUT = "8B8F93" # 深色页弱化
DARKRULE = "3A3E42"

F_HEI = "Microsoft YaHei"   # 标题/正文
F_MONO = "Consolas"         # 刊号/注记/数字

ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def rgb(h):
    return RGBColor.from_string(h)


def est_w(text, size, pad=0.30):
    units = sum(1.0 if ord(c) > 0x2000 else 0.55 for c in text)
    return units * size / 72.0 + pad


# ---------- 文本 ----------
def _set_run_font(run, name=F_HEI, size=12, bold=False, color=INK, spc=None, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    rPr = run._r.get_or_add_rPr()
    if spc is not None:
        rPr.set("spc", str(int(spc)))
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)


def _apply_bullet(p, color=RED, char="▪", marL=0.2):
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(marL))))
    pPr.set("indent", "-" + str(int(Inches(marL))))
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = pPr.makeelement(qn("a:srgbClr"), {"val": color})
    buClr.append(srgb)
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(buClr)
    pPr.append(buFont)
    pPr.append(buChar)


def _fill_tf(tf, paras, anchor="t", align="l", wrap=True):
    tf.word_wrap = wrap
    tf.vertical_anchor = ANCH[anchor]
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    # 中文排版禁则：避头尾 + 标点悬挂
    try:
        bodyPr = tf._txBody.bodyPr
        bodyPr.set("eaLnBrk", "1")
        bodyPr.set("hangingPunct", "1")
    except Exception:
        pass
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        opts = para.get("o", {}) if isinstance(para, dict) else {}
        runs = para["r"] if isinstance(para, dict) else para
        p.alignment = ALIGN[opts.get("align", align)]
        if "sa" in opts:
            p.space_after = Pt(opts["sa"])
        if "sb" in opts:
            p.space_before = Pt(opts["sb"])
        if "line" in opts:
            p.line_spacing = opts["line"]
        if opts.get("bullet"):
            bc = opts["bullet"]
            _apply_bullet(p, color=bc[1] if isinstance(bc, (list, tuple)) else RED,
                          char=bc[0] if isinstance(bc, (list, tuple)) else "▪",
                          marL=opts.get("marL", 0.2))
        for run in runs:
            r = p.add_run()
            r.text = run[0]
            _set_run_font(r, name=run[4] if len(run) > 4 else F_HEI,
                          size=run[1], bold=(run[2] if len(run) > 2 else False),
                          color=(run[3] if len(run) > 3 else INK),
                          spc=(run[5] if len(run) > 5 else None))


def text(slide, x, y, w, h, paras, anchor="t", align="l", wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _fill_tf(box.text_frame, paras, anchor=anchor, align=align, wrap=wrap)
    return box


def stext(sh, paras, anchor="m", align="l", wrap=True, ml=0.08, mr=0.08, mt=0.02, mb=0.02):
    tf = sh.text_frame
    tf.margin_left, tf.margin_right = Inches(ml), Inches(mr)
    tf.margin_top, tf.margin_bottom = Inches(mt), Inches(mb)
    _fill_tf(tf, paras, anchor=anchor, align=align, wrap=wrap)
    return sh


# ---------- 基础形状 ----------
def rect(slide, x, y, w, h, fill=WHITE, line_c=None, line_w=1.0, round_=None):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if round_ is not None else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    if round_ is not None:
        try:
            sh.adjustments[0] = round_
        except Exception:
            pass
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    if line_c:
        sh.line.color.rgb = rgb(line_c)
        sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def fill_alpha(sh, transparency_pct):
    xClr = sh.fill.fore_color._color._xClr
    a = xClr.makeelement(qn("a:alpha"), {"val": str(int((100 - transparency_pct) * 1000))})
    xClr.append(a)


def _line(slide, x1, y1, x2, y2, color=RULE, weight=0.75):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def rule(slide, x, y, w, color=RULE, weight=0.75):
    return _line(slide, x, y, x + w, y, color, weight)


def vline(slide, x, y1, y2, color=RULE, weight=0.75):
    return _line(slide, x, y1, x, y2, color, weight)


def arrow_line(slide, x1, y1, x2, y2, color=RED, weight=1.2):
    """带箭头的细线（编辑风流程箭头）"""
    ln = _line(slide, x1, y1, x2, y2, color, weight)
    lnEl = ln.line._get_or_add_ln()
    tail = lnEl.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    lnEl.append(tail)
    return ln


def dot(slide, cx, cy, d, color=RED):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2),
                                Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(color)
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def sq(slide, x, y, s, color=RED):
    sh = rect(slide, x, y, s, s, fill=color)
    return sh


def chip(slide, x, y, w, h, label, fill="F1EEE6", color=INK, size=10.5, bold=True, line_c=None):
    """极少使用的胶囊（编辑风基本不用卡片，此函数保留备用）"""
    sh = rect(slide, x, y, w, h, fill=fill, line_c=line_c, round_=0.5)
    stext(sh, [[(label, size, bold, color)]], anchor="m", align="c", ml=0.04, mr=0.04)
    return sh


def chips_row(slide, x, y, labels, h=0.34, size=10.5, gap=0.14, **kw):
    cx = x
    for lb in labels:
        w = est_w(lb, size)
        chip(slide, cx, y, w, h, lb, size=size, **kw)
        cx += w + gap
    return cx - gap


# ---------- 页面框架 ----------
def ghost(slide, x, y, txt, size=110, color=GHOST, w=5.0, h=1.9):
    """超大浅色序号（编辑风层次元素）"""
    return text(slide, x, y, w, h, [[(txt, size, True, color, F_MONO)]])


def kicker(slide, x, y, txt, color=RED, size=10.5, w=8.0):
    """红色等宽小标（SECTION 02 · MESHING）"""
    return text(slide, x, y, w, 0.3, [[(txt, size, True, color, F_MONO, 150)]])


def page_frame(slide, ch_en, ch_cn, page, total=21, dark=False, cover=False):
    """页眉页脚：顶部细线 + 左右运行头；底部细线 + 页码"""
    if cover:
        return
    mut = DARKMUT if dark else FAINT
    rl = DARKRULE if dark else RULE
    tcol = DARKTXT if dark else INK
    # 顶部运行头
    text(slide, 0.9, 0.42, 7.0, 0.26,
         [[("ANSYS WORKBENCH — 有限元分析学习总结", 8.5, False, mut, F_MONO, 100)]])
    text(slide, 8.4, 0.42, 4.03, 0.26,
         [[(ch_en + "  ", 8.5, False, mut, F_MONO, 100), (ch_cn, 9, False, mut)]], align="r")
    rule(slide, 0.9, 0.78, 11.53, color=rl)
    # 底部
    rule(slide, 0.9, 7.0, 11.53, color=rl)
    text(slide, 0.9, 7.08, 6.0, 0.26,
         [[("基于 e1-2 ~ e1-10 课程学习笔记整理 · 2026-09", 8.5, False, mut, F_MONO, 80)]])
    text(slide, 10.43, 7.02, 2.0, 0.34,
         [[("%02d" % page, 14, True, tcol, F_MONO), (" / %d" % total, 9, False, mut, F_MONO)]],
         align="r")


def title_block(slide, kick, title, standfirst=None, page_w=9.5):
    """内容页标题区：红色 kicker + 大标题 + 导语"""
    kicker(slide, 0.9, 1.02, kick)
    text(slide, 0.9, 1.34, page_w, 0.78, [[(title, 32, True, INK)]])
    if standfirst:
        text(slide, 0.9, 2.18, page_w + 0.6, 0.42,
             [[(standfirst, 12.5, False, GRAY)]], align="l")


def bullet_list(slide, x, y, w, items, size=12, gap=9, line=1.3, color=RED, tcolor=INK,
                bold_first=False, h=4.0):
    """方块悬挂列表；items: str 或 (term, desc)"""
    paras = []
    for it in items:
        if isinstance(it, tuple):
            t1, t2 = it
            runs = [(t1, size, True, tcolor)]
            if t2:
                runs.append(("　" + t2, size, False, GRAY))
            paras.append({"r": runs, "o": {"bullet": ("▪", color), "line": line, "sa": gap}})
        else:
            paras.append({"r": [(it, size, False, tcolor)],
                          "o": {"bullet": ("▪", color), "line": line, "sa": gap}})
    return text(slide, x, y, w, h, paras)


def label_rows(slide, x, y, w, rows, gap=0.14, size=11.5, label_color=INK,
               text_color=GRAY, line=1.3, label_gap=0.22):
    """「标签 | 说明」定义列表：标签列固定宽，说明列对齐；rows: [(label, desc), ...]"""
    import math
    lw = max(est_w(lb, size, pad=0.08) for lb, _ in rows)
    cy = y
    for lb, desc in rows:
        dw = w - lw - label_gap
        chars = max(1, int(dw / (size / 72.0)))
        units = sum(1.0 if ord(c) > 0x2000 else 0.55 for c in desc)
        nlines = max(1, math.ceil(units / chars))
        h = nlines * (size / 72.0) * line + 0.08
        text(slide, x, cy + 0.01, lw, 0.3, [[(lb, size, True, label_color)]])
        text(slide, x + lw + label_gap, cy, dw, h,
             [{"r": [(desc, size, False, text_color)], "o": {"line": line}}])
        cy += max(h, 0.26) + gap
    return cy
